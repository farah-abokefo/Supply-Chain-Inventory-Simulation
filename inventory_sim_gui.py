"""
Inventory Simulation GUI
- English UI
- Light / Dark theme option
- Run simulation, view plot, view results
- Save Excel, PDF, PowerPoint reports

Save as: inventory_sim_gui.py
Run: python inventory_sim_gui.py
"""

import sys
import random
import math
import tkinter as tk
from tkinter import ttk, messagebox, filedialog

# plotting
import matplotlib
matplotlib.use("TkAgg")
from matplotlib.figure import Figure
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg

# optional report libraries
try:
    from openpyxl import Workbook
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from pptx import Presentation
except Exception:
    # We'll handle missing libs at save time with an informative error
    pass


# ----------------------------
# Simulation core (class)
# ----------------------------
class InventorySimulation:
    def __init__(self, days=100, initial_inventory=120, reorder_point=40,
                 order_quantity=90, lead_time=5, holding_cost=0.6,
                 shortage_cost=5, ordering_cost=60, demand_min=5, demand_max=20):
        self.days = int(days)
        self.initial_inventory = int(initial_inventory)
        self.reorder_point = int(reorder_point)
        self.order_quantity = float(order_quantity)
        self.lead_time = int(lead_time)
        self.holding_cost = float(holding_cost)
        self.shortage_cost = float(shortage_cost)
        self.ordering_cost = float(ordering_cost)
        self.demand_min = int(demand_min)
        self.demand_max = int(demand_max)

        # results
        self.inventory_levels = []
        self.total_holding_cost = 0.0
        self.total_shortage_cost = 0.0
        self.total_ordering_cost = 0.0
        self.stockout_days = 0
        self.final_inventory = self.initial_inventory

    def generate_daily_demand(self):
        return random.randint(self.demand_min, self.demand_max)

    def compute_eoq(self, annual_demand=3000):
        # approximate H as daily holding * 365
        H = self.holding_cost * 365
        S = self.ordering_cost
        D = annual_demand
        if H <= 0:
            return None
        return math.sqrt((2 * D * S) / H)

    def run(self):
        self.inventory_levels = []
        self.total_holding_cost = 0.0
        self.total_shortage_cost = 0.0
        self.total_ordering_cost = 0.0
        self.stockout_days = 0
        inventory = self.initial_inventory
        pending_orders = []

        for day in range(1, self.days + 1):
            # receive orders
            for order in pending_orders[:]:
                if order[0] == day:
                    inventory += order[1]
                    pending_orders.remove(order)

            demand = self.generate_daily_demand()
            if inventory >= demand:
                inventory -= demand
            else:
                shortage = demand - inventory
                self.total_shortage_cost += shortage * self.shortage_cost
                inventory = 0
                self.stockout_days += 1

            self.total_holding_cost += inventory * self.holding_cost

            if inventory <= self.reorder_point:
                arrival = day + self.lead_time
                pending_orders.append((arrival, int(self.order_quantity)))
                self.total_ordering_cost += self.ordering_cost

            self.inventory_levels.append(inventory)

        self.final_inventory = inventory

    def summary(self):
        total_cost = self.total_holding_cost + self.total_shortage_cost + self.total_ordering_cost
        return {
            "days": self.days,
            "final_inventory": self.final_inventory,
            "holding_cost": round(self.total_holding_cost, 2),
            "shortage_cost": round(self.total_shortage_cost, 2),
            "ordering_cost": round(self.total_ordering_cost, 2),
            "total_cost": round(total_cost, 2),
            "stockout_days": self.stockout_days,
            "inventory_levels": self.inventory_levels
        }


# ----------------------------
# GUI Application
# ----------------------------
class InventoryApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("Supply Chain Inventory Simulation")
        self.geometry("980x640")
        self.resizable(True, True)

        # theme state
        self.theme = tk.StringVar(value="Light")

        # create UI
        self.create_widgets()
        self.apply_theme()

        # simulation object placeholder
        self.sim = None

    def create_widgets(self):
        # Left frame: inputs and controls
        left = ttk.Frame(self)
        left.pack(side=tk.LEFT, fill=tk.Y, padx=10, pady=10)

        # Parameters
        params = [
            ("Simulation Days", "100"),
            ("Initial Inventory", "120"),
            ("Reorder Point (ROP)", "40"),
            ("Order Quantity", "90"),
            ("Lead Time (days)", "5"),
            ("Holding Cost (per unit/day)", "0.6"),
            ("Shortage Cost (per unit)", "5"),
            ("Ordering Cost (per order)", "60"),
            ("Demand Min (per day)", "5"),
            ("Demand Max (per day)", "20")
        ]
        self.entries = {}
        row = 0
        for label_text, default in params:
            lbl = ttk.Label(left, text=label_text)
            lbl.grid(row=row, column=0, sticky=tk.W, pady=4)
            ent = ttk.Entry(left, width=14)
            ent.grid(row=row, column=1, pady=4, padx=6)
            ent.insert(0, default)
            self.entries[label_text] = ent
            row += 1

        # Theme toggle
        theme_frame = ttk.LabelFrame(left, text="Theme")
        theme_frame.grid(row=row, column=0, columnspan=2, pady=8, sticky=tk.EW)
        ttk.Radiobutton(theme_frame, text="Light", variable=self.theme, value="Light", command=self.apply_theme).pack(side=tk.LEFT, padx=6)
        ttk.Radiobutton(theme_frame, text="Dark", variable=self.theme, value="Dark", command=self.apply_theme).pack(side=tk.LEFT, padx=6)
        row += 1

        # Buttons
        btn_run = ttk.Button(left, text="Run Simulation", command=self.on_run)
        btn_run.grid(row=row, column=0, pady=6, sticky=tk.EW)
        btn_clear = ttk.Button(left, text="Clear Plot", command=self.clear_plot)
        btn_clear.grid(row=row, column=1, pady=6, sticky=tk.EW)
        row += 1

        btn_save_excel = ttk.Button(left, text="Save Excel", command=self.save_excel)
        btn_save_excel.grid(row=row, column=0, pady=6, sticky=tk.EW)
        btn_save_pdf = ttk.Button(left, text="Save PDF", command=self.save_pdf)
        btn_save_pdf.grid(row=row, column=1, pady=6, sticky=tk.EW)
        row += 1

        btn_save_ppt = ttk.Button(left, text="Save PowerPoint", command=self.save_ppt)
        btn_save_ppt.grid(row=row, column=0, columnspan=2, pady=6, sticky=tk.EW)
        row += 1

        # Results box
        res_frame = ttk.LabelFrame(left, text="Results")
        res_frame.grid(row=row, column=0, columnspan=2, pady=10, sticky=tk.EW)
        self.results_text = tk.Text(res_frame, height=10, width=34, wrap=tk.WORD)
        self.results_text.pack(fill=tk.BOTH, expand=True)
        self.results_text.insert(tk.END, "Run the simulation to see results here.")
        self.results_text.config(state=tk.DISABLED)

        # Right frame: plot
        right = ttk.Frame(self)
        right.pack(side=tk.RIGHT, fill=tk.BOTH, expand=True, padx=10, pady=10)

        fig = Figure(figsize=(6, 4), dpi=100)
        self.ax = fig.add_subplot(111)
        self.ax.set_title("Inventory Level Over Time")
        self.ax.set_xlabel("Day")
        self.ax.set_ylabel("Inventory Level")
        self.canvas = FigureCanvasTkAgg(fig, master=right)
        self.canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)

    def apply_theme(self):
        t = self.theme.get()
        if t == "Dark":
            bg = "#2b2b2b"
            fg = "#eaeaea"
            entry_bg = "#3c3f41"
        else:
            bg = "#f5f5f5"
            fg = "#000000"
            entry_bg = "#ffffff"

        style = ttk.Style(self)
        # general
        style.configure("TFrame", background=bg)
        style.configure("TLabel", background=bg, foreground=fg)
        style.configure("TLabelFrame", background=bg, foreground=fg)
        style.configure("TButton", background=bg)
        style.configure("TEntry", fieldbackground=entry_bg, background=entry_bg, foreground=fg)
        style.configure("TRadiobutton", background=bg, foreground=fg)

        self.configure(background=bg)
        # update text widget colors
        self.results_text.config(bg=entry_bg, fg=fg, insertbackground=fg)

    def get_parameters(self):
        try:
            params = {
                "days": int(self.entries["Simulation Days"].get()),
                "initial_inventory": int(self.entries["Initial Inventory"].get()),
                "reorder_point": int(self.entries["Reorder Point (ROP)"].get()),
                "order_quantity": float(self.entries["Order Quantity"].get()),
                "lead_time": int(self.entries["Lead Time (days)"].get()),
                "holding_cost": float(self.entries["Holding Cost (per unit/day)"].get()),
                "shortage_cost": float(self.entries["Shortage Cost (per unit)"].get()),
                "ordering_cost": float(self.entries["Ordering Cost (per order)"].get()),
                "dmin": int(self.entries["Demand Min (per day)"].get()),
                "dmax": int(self.entries["Demand Max (per day)"].get())
            }
            if params["dmin"] > params["dmax"]:
                raise ValueError("Demand min must be <= demand max")
            return params
        except Exception as e:
            messagebox.showerror("Invalid parameter", f"Please check parameters.\n\n{e}")
            return None

    def on_run(self):
        params = self.get_parameters()
        if not params:
            return
        self.sim = InventorySimulation(
            days=params["days"],
            initial_inventory=params["initial_inventory"],
            reorder_point=params["reorder_point"],
            order_quantity=params["order_quantity"],
            lead_time=params["lead_time"],
            holding_cost=params["holding_cost"],
            shortage_cost=params["shortage_cost"],
            ordering_cost=params["ordering_cost"],
            demand_min=params["dmin"],
            demand_max=params["dmax"],
        )
        self.sim.run()
        summary = self.sim.summary()
        self.show_results(summary)
        self.plot_inventory(summary["inventory_levels"])

    def show_results(self, summary):
        self.results_text.config(state=tk.NORMAL)
        self.results_text.delete("1.0", tk.END)
        lines = [
            f"Simulation Days: {summary['days']}",
            f"Final Inventory Level: {summary['final_inventory']}",
            f"Holding Cost: {summary['holding_cost']}",
            f"Shortage Cost: {summary['shortage_cost']}",
            f"Ordering Cost: {summary['ordering_cost']}",
            f"TOTAL SYSTEM COST: {summary['total_cost']}",
            f"Stockout Days: {summary['stockout_days']}"
        ]
        self.results_text.insert(tk.END, "\n".join(lines))
        self.results_text.config(state=tk.DISABLED)

    def plot_inventory(self, levels):
        self.ax.clear()
        self.ax.plot(range(1, len(levels) + 1), levels, marker="", linewidth=1)
        self.ax.set_title("Inventory Level Over Time")
        self.ax.set_xlabel("Day")
        self.ax.set_ylabel("Inventory Level")
        self.ax.grid(True)
        self.canvas.draw()

    def clear_plot(self):
        self.ax.clear()
        self.ax.set_title("Inventory Level Over Time")
        self.ax.set_xlabel("Day")
        self.ax.set_ylabel("Inventory Level")
        self.canvas.draw()
        self.results_text.config(state=tk.NORMAL)
        self.results_text.delete("1.0", tk.END)
        self.results_text.insert(tk.END, "Run the simulation to see results here.")
        self.results_text.config(state=tk.DISABLED)

    # -------------------------
    # Save helpers
    # -------------------------
    def save_excel(self):
        if not self.sim:
            messagebox.showinfo("No data", "Run the simulation first.")
            return
        try:
            path = filedialog.asksaveasfilename(defaultextension=".xlsx",
                                                filetypes=[("Excel files", "*.xlsx")])
            if not path:
                return
            wb = Workbook()
            ws = wb.active
            ws.title = "Inventory Results"
            ws.append(["Day", "Inventory Level"])
            for i, v in enumerate(self.sim.inventory_levels, start=1):
                ws.append([i, v])
            wb.save(path)
            messagebox.showinfo("Saved", f"Excel saved to:\n{path}")
        except Exception as e:
            messagebox.showerror("Error saving Excel", str(e))

    def save_pdf(self):
        if not self.sim:
            messagebox.showinfo("No data", "Run the simulation first.")
            return
        try:
            path = filedialog.asksaveasfilename(defaultextension=".pdf",
                                                filetypes=[("PDF files", "*.pdf")])
            if not path:
                return
            # Check if reportlab is available
            try:
                pdf = SimpleDocTemplate(path, pagesize=A4)
                styles = getSampleStyleSheet()
            except Exception as e:
                raise RuntimeError("reportlab is required for PDF export. Install with: pip install reportlab") from e

            elements = []
            elements.append(Paragraph("Supply Chain Inventory Simulation Report", styles["Title"]))
            elements.append(Spacer(1, 12))
            eoq = self.sim.compute_eoq()
            elements.append(Paragraph(f"EOQ (approx): {round(eoq,2) if eoq else 'N/A'}", styles["Normal"]))
            s = self.sim.summary()
            elements.append(Paragraph(f"Simulation Days: {s['days']}", styles["Normal"]))
            elements.append(Paragraph(f"Final Inventory Level: {s['final_inventory']}", styles["Normal"]))
            elements.append(Paragraph(f"Holding Cost: {s['holding_cost']}", styles["Normal"]))
            elements.append(Paragraph(f"Shortage Cost: {s['shortage_cost']}", styles["Normal"]))
            elements.append(Paragraph(f"Ordering Cost: {s['ordering_cost']}", styles["Normal"]))
            elements.append(Paragraph(f"Total System Cost: {s['total_cost']}", styles["Normal"]))
            elements.append(Paragraph(f"Stockout Days: {s['stockout_days']}", styles["Normal"]))
            pdf.build(elements)
            messagebox.showinfo("Saved", f"PDF saved to:\n{path}")
        except Exception as e:
            messagebox.showerror("Error saving PDF", str(e))

    def save_ppt(self):
        if not self.sim:
            messagebox.showinfo("No data", "Run the simulation first.")
            return
        try:
            path = filedialog.asksaveasfilename(defaultextension=".pptx",
                                                filetypes=[("PowerPoint", "*.pptx")])
            if not path:
                return
            # Check python-pptx
            try:
                prs = Presentation()
            except Exception as e:
                raise RuntimeError("python-pptx is required for PPT export. Install with: pip install python-pptx") from e

            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            s = self.sim.summary()
            eoq = self.sim.compute_eoq()
            title.text = "Inventory Simulation Results"
            content.text = (f"EOQ ≈ {round(eoq,2) if eoq else 'N/A'}\n"
                            f"Total Cost: {s['total_cost']}\n"
                            f"Holding Cost: {s['holding_cost']}\n"
                            f"Shortage Cost: {s['shortage_cost']}\n"
                            f"Ordering Cost: {s['ordering_cost']}\n"
                            f"Stockout Days: {s['stockout_days']}")
            prs.save(path)
            messagebox.showinfo("Saved", f"PPT saved to:\n{path}")
        except Exception as e:
            messagebox.showerror("Error saving PPT", str(e))


if __name__ == "__main__":
    app = InventoryApp()
    app.mainloop()
